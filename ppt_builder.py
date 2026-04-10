# from pptx import Presentation
# from pptx.util import Inches, Pt
# from pptx.dml.color import RGBColor
# from pptx.enum.text import PP_ALIGN

# DARK_BG   = RGBColor(0x1A, 0x1A, 0x2E)
# BLUE_BG   = RGBColor(0x16, 0x21, 0x3E)
# ACCENT    = RGBColor(0x0F, 0x3A, 0x60)
# WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
# LIGHT_BLUE= RGBColor(0xAA, 0xCC, 0xEE)
# BRIGHT    = RGBColor(0x44, 0xAA, 0xFF)

# def add_textbox(slide, text, left, top, width, height,
#                 size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
#     box = slide.shapes.add_textbox(
#         Inches(left), Inches(top), Inches(width), Inches(height)
#     )
#     tf = box.text_frame
#     tf.word_wrap = True
#     p = tf.paragraphs[0]
#     p.alignment = align
#     run = p.add_run()
#     run.text = text
#     run.font.size = Pt(size)
#     run.font.bold = bold
#     run.font.color.rgb = color

# def set_bg(slide, color):
#     slide.background.fill.solid()
#     slide.background.fill.fore_color.rgb = color

# def make_title_slide(prs, title, subtitle):
#     slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
#     set_bg(slide, DARK_BG)

#     # Left accent bar
#     bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.2), Inches(7.5))
#     bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()

#     add_textbox(slide, title,       0.5, 2.0, 12.0, 1.5, size=44, bold=True)
#     add_textbox(slide, subtitle,    0.5, 3.7, 10.0, 1.0, size=22, color=LIGHT_BLUE)
#     add_textbox(slide, "Made with AI", 0.5, 6.5, 6.0, 0.5, size=12, color=LIGHT_BLUE)

# def make_content_slide(prs, heading, bullets, speaker_note=""):
#     slide = prs.slides.add_slide(prs.slide_layouts[6])
#     set_bg(slide, BLUE_BG)

#     # Top strip
#     strip = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.08))
#     strip.fill.solid(); strip.fill.fore_color.rgb = ACCENT; strip.line.fill.background()

#     # Heading
#     add_textbox(slide, heading, 0.4, 0.2, 12.0, 0.9, size=30, bold=True)

#     # Bullets
#     for i, bullet in enumerate(bullets):
#         top = 1.5 + i * 0.85

#         # Background box per bullet
#         bg = slide.shapes.add_shape(1, Inches(0.4), Inches(top), Inches(12.3), Inches(0.65))
#         bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x33)
#         bg.line.color.rgb = RGBColor(0x22, 0x55, 0x88); bg.line.width = Pt(0.5)

#         add_textbox(slide, "›", 0.5, top+0.05, 0.3, 0.55, size=18, bold=True, color=BRIGHT)
#         add_textbox(slide, bullet, 0.85, top+0.05, 11.7, 0.55, size=17)

#     if speaker_note:
#         slide.notes_slide.notes_text_frame.text = speaker_note

# def build_ppt(data: dict, filename: str = "output.pptx"):
#     prs = Presentation()
#     prs.slide_width  = Inches(13.33)
#     prs.slide_height = Inches(7.5)

#     make_title_slide(prs, data["title"], data["subtitle"])
#     print(f"Title slide: {data['title']}")

#     for i, s in enumerate(data["slides"]):
#         make_content_slide(prs, s["heading"], s["bullets"], s.get("speaker_note",""))
#         print(f"Slide {i+1}: {s['heading']}")

#     prs.save(filename)
#     print(f"\nSaved: {filename}")

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Themes
THEMES = {
    'navy':   (RGBColor(0x1A,0x1A,0x2E), RGBColor(0x16,0x21,0x3E), RGBColor(0x0F,0x3A,0x60)),
    'dark':   (RGBColor(0x1C,0x1C,0x1C), RGBColor(0x2C,0x2C,0x2C), RGBColor(0x44,0x44,0x44)),
    'green':  (RGBColor(0x1A,0x2E,0x1A), RGBColor(0x16,0x3E,0x21), RGBColor(0x0F,0x60,0x3A)),
    'purple': (RGBColor(0x1A,0x0A,0x2E), RGBColor(0x2A,0x10,0x4E), RGBColor(0x3A,0x15,0x6E)),
}

WHITE = RGBColor(255,255,255)

def build_ppt(data, filename="output.pptx", theme="navy"):
    DARK_BG, BLUE_BG, ACCENT = THEMES.get(theme, THEMES['navy'])

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK_BG

    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(10), Inches(2))
    title_box.text_frame.text = data["title"]

    # Slides
    for s in data["slides"]:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = BLUE_BG

        heading = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(1))
        heading.text_frame.text = s["heading"]

        content = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(10), Inches(4))
        tf = content.text_frame

        for b in s["bullets"]:
            p = tf.add_paragraph()
            p.text = b
            p.level = 0

    prs.save(filename)